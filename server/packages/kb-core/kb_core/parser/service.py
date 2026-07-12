"""Document parser module: dispatches to MinerU (PDF/images), LlamaIndex readers (office/text),
or stdlib email parser based on file type."""

from __future__ import annotations

import logging
import re as re_mod
import tempfile

logger = logging.getLogger("kb_core.parser")
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

SUPPORTED_TYPES = {
    "pdf",
    "docx",
    "xlsx",
    "pptx",
    "md",
    "txt",
    "csv",
    "jpg",
    "jpeg",
    "png",
    "bmp",
    "tiff",
    "tif",
    "msg",
    "eml",
}


@dataclass
class ParseResult:
    text: str
    metadata: dict = field(default_factory=dict)


def _get_file_type(file_name: str) -> str:
    suffix = Path(file_name).suffix.lower().lstrip(".")
    return suffix


def _parse_pdf(file_path: str) -> str:
    """PDF 解析使用 PyMuPDF 结构化提取，保留段落、标题层次和页码。"""
    import fitz
    import re as _re

    doc = fitz.open(file_path)
    try:
        # 先扫一遍全文，获取字体大小分布，确定标题阈值
        all_sizes: list[float] = []
        pages_data: list[list[dict]] = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            if span.get("text", "").strip():
                                all_sizes.append(span.get("size", 0))
            pages_data.append(blocks)

        # 排序取最大字体、中位字体
        all_sizes.sort(reverse=True)
        max_size = all_sizes[0] if all_sizes else 12
        median_size = all_sizes[len(all_sizes) // 2] if all_sizes else 11
        heading_threshold = max(median_size * 1.3, 14)

        pages = []
        for page_num in range(len(doc)):
            blocks = pages_data[page_num]
            page_blocks: list[tuple[str, float, float]] = []  # (text, y, x)

            for block in blocks:
                if block.get("type") == 0:
                    block_text_parts: list[str] = []
                    prev_is_heading = False
                    for line in block.get("lines", []):
                        spans = [s["text"] for s in line.get("spans", [])]
                        line_text = "".join(spans).strip()
                        if not line_text:
                            continue

                        span = line.get("spans", [{}])[0]
                        font_size = span.get("size", 0)
                        flags = span.get("flags", 0)
                        is_bold = bool(flags & 2)

                        if font_size >= heading_threshold + 4 or (is_bold and font_size >= heading_threshold + 2):
                            level = 1
                        elif font_size >= heading_threshold or is_bold:
                            level = 2
                        elif font_size > median_size * 1.1:
                            level = 3
                        else:
                            level = 0

                        if level > 0 and _re.match(r"^[▪•\-\*▶]", line_text.strip()):
                            level = 0

                        if level > 0:
                            text = line_text.strip()
                            text = text.replace("#", "")
                            num_match = _re.match(r"^\d+(\.\d+)+", text)
                            if num_match:
                                dots = num_match.group(0).count(".")
                                level = min(level + (2 if dots >= 2 else 1), 6)
                            line_text = f"{'#' * level} {text}"

                            if block_text_parts and not prev_is_heading:
                                block_text_parts.append("\n" + line_text)
                            else:
                                block_text_parts.append(line_text)
                            prev_is_heading = True
                        else:
                            if block_text_parts and not prev_is_heading and not block_text_parts[-1].endswith("\n"):
                                block_text_parts[-1] += " " + line_text
                            else:
                                block_text_parts.append(line_text)
                            prev_is_heading = False

                    if block_text_parts:
                        text = "".join(block_text_parts)
                        bbox = block.get("bbox", (0, 0, 0, 0))
                        page_blocks.append((text, bbox[1], bbox[0]))  # (text, y, x)

            if page_blocks:
                pages.append(page_blocks)  # now list of (text, y, x) tuples

        # 按页码组装文本，检测表格并格式化行
        result_parts = []
        for page_idx, page_blocks in enumerate(pages):
            page_text = _assemble_page(page_blocks, 80)
            result_parts.append(f"\[[page:{page_idx+1}]\]\n{page_text}")

        result = "\n\n".join(result_parts)
        for zwc in ["​", "﻿", "­"]:
            result = result.replace(zwc, "")
        result = _re.sub(r"\n{3,}", "\n\n", result)
        # 列表项换行：i. ii. iii. a. b. c. 等序号单独成行
        result = _re.sub(r"([；;。])\s*([a-z]+\. )", r"\1\n\2", result)
        result = _re.sub(r" +", " ", result)
        return result
    finally:
        doc.close()


def _merge_table_cells(text: str, max_line: int = 60, min_rows: int = 3) -> str:
    """合并疑似表格的连续短行段落，用 | 分隔，再清理括号内误分割."""
    paragraphs = text.split("\n\n")
    result: list[str] = []
    i = 0
    while i < len(paragraphs):
        j = i
        shorts: list[str] = []
        while j < len(paragraphs):
            p = paragraphs[j].strip()
            if p.startswith("#") or not p:
                break
            if len(p) > max_line:
                break
            shorts.append(p)
            j += 1
        if len(shorts) >= min_rows:
            # 首行（表头）可能是空格分隔的多列：数据分级 存储要求 迁移约束
            parts = shorts[0].split(" ")
            if len(parts) >= 3 and all(len(p) < 15 for p in parts):
                shorts[0] = " | ".join(parts)
            merged = " | ".join(shorts)
            # 去掉括号内的 |（银行 | 卡→银行 卡）
            merged = re_mod.sub(r"（[^）]*?\|[^）]*?）", lambda m: m.group(0).replace(" | ", " "), merged)
            result.append(merged)
            i = j
        else:
            result.append(paragraphs[i])
            i += 1
    return "\n\n".join(result)


def _assemble_page(blocks: list[tuple[str, float, float]], short_threshold: int = 60) -> str:
    """将带位置信息的 block 组装成文本，检测表格并按行格式化."""
    if not blocks:
        return ""

    is_short = [len(t) <= short_threshold and not t.startswith("#") for t, _, _ in blocks]

    lines: list[str] = []
    i = 0
    while i < len(blocks):
        if is_short[i] and sum(is_short[i:i+3]) >= 3:
            # 收集表格区域块
            table_blocks: list[tuple[str, float, float]] = []
            while i < len(blocks) and is_short[i]:
                table_blocks.append(blocks[i])
                i += 1

            # 按 Y 分组（Y 差 < 20 同一行），同 X 合并
            y_rows: list[tuple[float, list[tuple[str, float]]]] = []  # (ref_y, [(text, x), ...])
            for text, y, x in table_blocks:
                added = False
                for ri, (ref_y, cells) in enumerate(y_rows):
                    if abs(ref_y - y) < 20:
                        for ci, (ct, cx) in enumerate(cells):
                            if abs(cx - x) < 30:
                                cells[ci] = (ct + text, cx)  # merge
                                added = True
                                break
                        if not added:
                            cells.append((text, x))
                        added = True
                        break
                if not added:
                    y_rows.append((y, [(text, x)]))

            # 跨行合并续行：非空单元格匹配前一行同 X 的列
            changed = True
            while changed:
                changed = False
                for ri in range(len(y_rows) - 1, 0, -1):
                    cur = y_rows[ri]
                    prev = y_rows[ri - 1]
                    _zw_chars = " \t\n\r ​﻿‌‍"
                    non_empty = [(t, x) for t, x in cur[1] if t.strip(_zw_chars)]
                    if len(non_empty) == 1 and len(prev[1]) > 0:
                        ct, cx = non_empty[0]
                        for pi, (pt, px) in enumerate(prev[1]):
                            if abs(px - cx) < 50:
                                prev[1][pi] = (pt + ct, px)
                                del y_rows[ri]
                                changed = True
                                break

            # 首行如果是空格分隔的短词（表头一行多列），拆开
            if y_rows:
                first_cells = y_rows[0][1]
                if len(first_cells) == 1:
                    text = first_cells[0][0]
                    parts = text.split(" ")
                    if len(parts) >= 3 and all(len(p) < 15 for p in parts):
                        y_rows[0] = (y_rows[0][0], [(p, 0) for p in parts])

            for _, cells in y_rows:
                cells.sort(key=lambda t: t[1])
                # 去掉空单元格
                def _has_content(t: str) -> bool:
                    import unicodedata
                    cleaned = ''.join(c for c in t if unicodedata.category(c) not in ('Cc', 'Cf', 'Zs'))
                    return bool(cleaned.strip())
                non_empty = [t for t, _ in cells if _has_content(t)]
                if not non_empty:
                    continue
                line = " | ".join(non_empty)
                line = re_mod.sub(r"（[^）]*?\|[^）]*?）", lambda m: m.group(0).replace(" | ", " "), line)
                lines.append(line)
        else:
            text = blocks[i][0]
            if text:
                lines.append(text)
            i += 1

    return "\n\n".join(lines)


def _parse_with_llama_index(file_type: str, file_path: str) -> str:
    """非 PDF/图片格式解析。使用轻量库替代 llama-index-readers-file。"""
    try:
        if file_type == "docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            lines = []
            for p in doc.paragraphs:
                t = p.text.strip()
                if not t:
                    lines.append("")
                    continue
                style = p.style.name.lower() if p.style else ""
                if style.startswith("heading 1"):
                    lines.append(f"# {t}")
                elif style.startswith("heading 2"):
                    lines.append(f"## {t}")
                elif style.startswith("heading 3"):
                    lines.append(f"### {t}")
                elif style.startswith("heading"):
                    # Heading 4+ or custom heading styles
                    level = min(6, int(style.split()[-1]) if style.split()[-1].isdigit() else 4)
                    lines.append(f"{'#' * level} {t}")
                else:
                    lines.append(t)
            return "\n".join(lines)
        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            lines = []
            for slide in prs.slides:
                # Title shape first (if exists)
                title_shape = slide.shapes.title
                if title_shape and title_shape.text.strip():
                    lines.append(f"# {title_shape.text.strip()}")
                # Other shapes (body)
                for shape in slide.shapes:
                    if shape == title_shape:
                        continue
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
                lines.append("")  # slide separator
            return "\n".join(lines)
        elif file_type in ("xlsx", "xls"):
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            rows = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text = " ".join(str(c) for c in row if c is not None)
                    if text.strip():
                        rows.append(text)
            return "\n".join(rows)
        elif file_type == "md":
            return Path(file_path).read_text(encoding="utf-8", errors="replace")
        elif file_type == "csv":
            import csv
            with open(file_path, newline="", encoding="utf-8", errors="replace") as f:
                return "\n".join(" ".join(row) for row in csv.reader(f))
        else:
            return Path(file_path).read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        logger.warning("Fallback parser failed for %s: %s", file_type, e)
        return Path(file_path).read_text(encoding="utf-8", errors="replace")


def _parse_image(file_path: str) -> str:
    """图片 OCR 使用 paddleocr 或 MinerU."""
    try:
        from magic_pdf.ocr import OCREngine
        ocr = OCREngine()
        result = ocr.ocr_file(file_path)
        return "\n".join(line["text"] for line in result)
    except ImportError:
        # fallback: 无 OCR
        return ""


def parse_document(
    content: bytes, file_name: str, file_type: Optional[str] = None
) -> ParseResult:
    """
    解析文档内容，返回文本和元数据。

    Args:
        content: 文件二进制内容
        file_name: 原始文件名（用于判断类型）
        file_type: 可选的文件类型覆盖

    Returns:
        ParseResult 包含解析后的文本和元数据
    """
    if file_type is None:
        file_type = _get_file_type(file_name)

    # 写临时文件供解析器读取
    suffix = Path(file_name).suffix
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        if file_type == "pdf":
            text = _parse_pdf(tmp_path)
        elif file_type in {"jpg", "jpeg", "png", "bmp", "tiff", "tif"}:
            text = _parse_image(tmp_path)
        elif file_type in {"msg", "eml"}:
            from email import policy
            from email.parser import BytesParser
            text = ""
            with open(tmp_path, "rb") as f:
                msg = BytesParser(policy=policy.default).parse(f)
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        text += part.get_content() or ""
            else:
                text = msg.get_content() or ""
        elif file_type in {"txt", "md", "csv", "json", "xml", "yaml", "yml", "ini", "cfg", "conf", "log", "sh", "py", "js", "ts", "html", "css"}:
            with open(tmp_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
        else:
            text = _parse_with_llama_index(file_type, tmp_path)
    finally:
        Path(tmp_path).unlink(missing_ok=True)

    return ParseResult(
        text=text.strip(),
        metadata={"file_name": file_name, "file_type": file_type, "char_count": len(text)},
    )


class ParserService:
    """文档解析服务的可注入包装."""

    def parse(self, content: bytes, file_name: str, file_type: Optional[str] = None) -> ParseResult:
        return parse_document(content, file_name, file_type)
