"""Office document text extraction and HTML conversion for local preview."""

from __future__ import annotations

import io
import logging
from typing import Optional

logger = logging.getLogger("kb_biz.services.preview")

OFFICE_EXTENSIONS = frozenset({".docx", ".pptx", ".ppt", ".xlsx", ".xls"})

_HTML_STYLES = """\
body {
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Noto Sans SC", sans-serif;
  font-size: 14px; line-height: 1.7; color: #1a1a1a;
  padding: 24px 32px; max-width: 100%; margin: 0 auto; background: #fff;
}
h1 { font-size: 22px; margin: 20px 0 10px; border-bottom: 2px solid #eee; padding-bottom: 6px; }
h2 { font-size: 18px; margin: 16px 0 8px; }
h3 { font-size: 15px; margin: 12px 0 6px; }
table { border-collapse: collapse; width: 100%; margin: 12px 0; font-size: 13px; }
table thead th { position: sticky; top: 0; z-index: 1; }
table td, table th { border: 1px solid #ccc; padding: 6px 10px; text-align: left; vertical-align: top; white-space: nowrap; }
table th { background: #f5f5f5; font-weight: 600; }
table tbody tr:nth-child(even) { background: #f9f9f9; }
table tbody tr:hover { background: #e8f4fd; }
.excel-sheet { overflow: auto; max-height: 65vh; border: 1px solid #ddd; border-radius: 6px; margin: 12px 0; }
img { max-width: 100%; height: auto; }
.slide { margin: 16px 0; padding: 16px; border: 1px solid #ddd; border-radius: 8px; background: #fafafa; }
.slide h3 { margin-top: 0; color: #555; }
@media (prefers-color-scheme: dark) {
  body { background: #1a1a2e; color: #e0e0e0; }
  table td, table th { border-color: #444; }
  table th { background: #2a2a3e; }
  .slide { background: #16162a; border-color: #333; }
  .slide h3 { color: #aaa; }
}"""


_DOWNLOAD_PLACEHOLDER = "__DOWNLOAD_URL__"


def _wrap_html(body: str, download_url: str = "") -> str:
    banner = ""
    if download_url:
        banner = f"""\
<div style="background:#fff3cd;color:#856404;padding:8px 16px;font-size:12px;border-bottom:1px solid #ffc107;text-align:center;flex-shrink:0">
预览效果可能与实际文档有差异，建议<a href="{download_url}" download style="color:#533f03;font-weight:600">下载</a>查看原文件
</div>"""
    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1.0">
<style>{_HTML_STYLES}</style></head>
<body style="display:flex;flex-direction:column;height:100vh;padding:0;overflow:hidden">
{banner}
<div style="flex:1;overflow:auto;padding:24px 32px">{body}</div>
</body></html>"""


def extract_text(file_name: str, data: bytes) -> Optional[str]:
    ext = _ext(file_name)
    if ext == ".docx":
        return _extract_docx(data)
    elif ext in (".xlsx", ".xls"):
        return _extract_xlsx(data)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(data)
    return None


def extract_html(file_name: str, data: bytes, download_url: str = "") -> Optional[str]:
    ext = _ext(file_name)
    if ext == ".docx":
        return _convert_docx_to_html(data, download_url)
    elif ext in (".xlsx", ".xls"):
        return _convert_xlsx_to_html(data, download_url)
    elif ext in (".pptx", ".ppt"):
        return _convert_pptx_to_html(data, download_url)
    return None


def _ext(name: str) -> str:
    return "." + name.rsplit(".", 1)[-1].lower() if "." in name else ""


# ── DOCX ──────────────────────────────────────────────

def _convert_docx_to_html(data: bytes, download_url: str = "") -> Optional[str]:
    try:
        import mammoth
        result = mammoth.convert_to_html(io.BytesIO(data))
        if not result.value.strip():
            return None
        return _wrap_html(result.value, download_url)
    except Exception as e:
        logger.warning("docx HTML failed: %s", e)
        return None


def _extract_docx(data: bytes) -> Optional[str]:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for i, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rt = " | ".join(cells)
                if rt.strip(" |"):
                    rows.append(rt)
            if rows:
                lines.append("")
                lines.append(f"[表格 {i}]")
                lines.extend(rows)
        return "\n".join(lines) if lines else None
    except Exception as e:
        logger.warning("docx text fallback failed: %s", e)
        return None


# ── XLSX / XLS ────────────────────────────────────────

def _convert_xlsx_to_html(data: bytes, download_url: str = "") -> Optional[str]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_rows = list(ws.iter_rows(values_only=True))
            # find last non-empty row
            data_rows = []
            for row in reversed(all_rows):
                if any(c is not None for c in row):
                    data_rows.insert(0, row)
                elif data_rows:
                    break  # stop at first fully empty gap from bottom
            if not data_rows:
                continue
            # first data row → thead, rest → tbody
            header = data_rows[0]
            body = data_rows[1:]
            thead = "<thead><tr>" + "".join(
                f"<th>{str(c) if c is not None else ''}</th>" for c in header
            ) + "</tr></thead>"
            tbody = ""
            if body:
                tbody = "<tbody>" + "".join(
                    "<tr>" + "".join(
                        f"<td>{str(c) if c is not None else ''}</td>" for c in row
                    ) + "</tr>" for row in body
                ) + "</tbody>"
            parts.append(f"<h2>{sheet_name}</h2><div class='excel-sheet'><table>{thead}{tbody}</table></div>")
        wb.close()
        return _wrap_html("".join(parts), download_url) if parts else None
    except Exception as e:
        logger.warning("xlsx HTML failed: %s", e)
        return None


def _extract_xlsx(data: bytes) -> Optional[str]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) if c is not None else "" for c in row]
                line = "\t".join(vals).strip()
                if line:
                    rows.append(line)
            if rows:
                lines.append(f"=== {sheet_name} ===")
                lines.extend(rows)
                lines.append("")
        wb.close()
        return "\n".join(lines) if lines else None
    except Exception as e:
        logger.warning("xlsx text failed: %s", e)
        return None


# ── PPTX / PPT ────────────────────────────────────────

def _convert_pptx_to_html(data: bytes, download_url: str = "") -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides_html = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                paras = "".join(f"<p>{t}</p>" for t in texts)
                slides_html.append(f'<div class="slide"><h3>第 {i} 页</h3>{paras}</div>')
        return _wrap_html("".join(slides_html), download_url) if slides_html else None
    except Exception as e:
        logger.warning("pptx HTML failed: %s", e)
        return None


def _extract_pptx(data: bytes) -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                lines.append(f"--- 第 {i} 页 ---")
                lines.extend(texts)
                lines.append("")
        return "\n".join(lines) if lines else None
    except Exception as e:
        logger.warning("pptx text failed: %s", e)
        return None
